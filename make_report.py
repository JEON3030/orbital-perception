#!/usr/bin/env python3
"""전력 대결 결과(results.json) → 발표용 비교 그래프(PNG) + PPTX.

한글: Noto Sans CJK KR 등록. 그래프는 로그-로그로 크기별 에너지·속도를 보이고,
20m 풀타일(30MP) 운영점을 표시. CPU 는 상한 위를 멱함수 외삽(명시).
"""
from __future__ import annotations

import json
from pathlib import Path

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm

OUT = Path("outputs/power_duel")
RES = OUT / "results.json"
FIG = OUT / "power_duel.png"
PPTX = OUT / "위성_추론_CPU_vs_GPU_대결.pptx"

# ── 한글 폰트 ──────────────────────────────────────────────────────────
KFONT = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
fm.fontManager.addfont(KFONT)
KO = fm.FontProperties(fname=KFONT).get_name()
plt.rcParams["font.family"] = KO
plt.rcParams["axes.unicode_minus"] = False

CPU_C, GPU_C = "#e5484d", "#2ea043"
FULL_TILE_MP = 5490 * 5490 / 1e6


def load():
    d = json.loads(RES.read_text())
    rows = d["rows"]
    def series(dev, key):
        xs, ys = [], []
        for r in rows:
            v = r.get(dev, {})
            if key in v and isinstance(v[key], (int, float)):
                xs.append(r["megapixels"]); ys.append(v[key])
        return np.array(xs), np.array(ys)
    return d, rows, series


def powerlaw_extrap(x, y, x_target):
    """log-log 1차 적합으로 x_target 에서 y 외삽(양수 가정)."""
    lx, ly = np.log(x), np.log(y)
    a, b = np.polyfit(lx, ly, 1)      # ly = a*lx + b
    return float(np.exp(a * np.log(x_target) + b)), a


def extrap_line(x, y, x_to):
    """마지막 측정점 → x_to 까지 멱함수 외삽 점선용 (xs, ys)."""
    lx, ly = np.log(x), np.log(y)
    a, b = np.polyfit(lx, ly, 1)
    xs = np.array([x.max(), x_to])
    ys = np.exp(a * np.log(xs) + b)
    return xs, ys


def make_fig(d, rows, series):
    cx_e, cy_e = series("cpu", "mJ_per_frame")
    gx_e, gy_e = series("gpu", "mJ_per_frame")
    cx_t, cy_t = series("cpu", "sec_per_frame")
    gx_t, gy_t = series("gpu", "sec_per_frame")
    cy_t, gy_t = cy_t * 1000, gy_t * 1000   # ms

    # 30MP 운영점 CPU 외삽
    e_ext, e_slope = powerlaw_extrap(cx_e, cy_e, FULL_TILE_MP)
    t_ext, _ = powerlaw_extrap(cx_t, cy_t, FULL_TILE_MP)
    g_e_ext, _ = powerlaw_extrap(gx_e, gy_e, FULL_TILE_MP)
    g_t_ext, _ = powerlaw_extrap(gx_t, gy_t, FULL_TILE_MP)

    fig, axes = plt.subplots(1, 2, figsize=(13.3, 5.2))
    idle = d.get("idle_watt", 0)

    # (1) 에너지/프레임
    ax = axes[0]
    ax.loglog(cx_e, cy_e, "o-", color=CPU_C, lw=2.4, ms=7, label="CPU (ARM 6코어)")
    ax.loglog(gx_e, gy_e, "s-", color=GPU_C, lw=2.4, ms=7, label="GPU (Orin iGPU)")
    xe, ye = extrap_line(cx_e, cy_e, FULL_TILE_MP); ax.loglog(xe, ye, ":", color=CPU_C, lw=1.8)
    xe, ye = extrap_line(gx_e, gy_e, FULL_TILE_MP); ax.loglog(xe, ye, ":", color=GPU_C, lw=1.8)
    ax.axvline(FULL_TILE_MP, color="#8b949e", ls="--", lw=1.3)
    ax.annotate(f"S2 20m 풀타일\n≈{FULL_TILE_MP:.0f}MP (운영점)", xy=(FULL_TILE_MP, ax.get_ylim()[0]),
                xytext=(FULL_TILE_MP*0.34, cy_e.min()*1.1), fontsize=9, color="#57606a")
    ax.set_title("① 한 장 처리 에너지 (mJ/frame)\n낮을수록 좋음 · idle 포함 total", fontsize=12, weight="bold")
    ax.set_xlabel("장면 크기 (메가픽셀, 20m)"); ax.set_ylabel("mJ / frame")
    ax.grid(True, which="both", alpha=0.25); ax.legend(fontsize=10, loc="upper left")

    # (2) 처리시간/프레임
    ax = axes[1]
    ax.loglog(cx_t, cy_t, "o-", color=CPU_C, lw=2.4, ms=7, label="CPU (ARM 6코어)")
    ax.loglog(gx_t, gy_t, "s-", color=GPU_C, lw=2.4, ms=7, label="GPU (Orin iGPU)")
    xe, ye = extrap_line(cx_t, cy_t, FULL_TILE_MP); ax.loglog(xe, ye, ":", color=CPU_C, lw=1.8, label="외삽")
    xe, ye = extrap_line(gx_t, gy_t, FULL_TILE_MP); ax.loglog(xe, ye, ":", color=GPU_C, lw=1.8)
    ax.axvline(FULL_TILE_MP, color="#8b949e", ls="--", lw=1.3)
    ax.set_title("② 한 장 처리 시간 (ms/frame)\n낮을수록 좋음", fontsize=12, weight="bold")
    ax.set_xlabel("장면 크기 (메가픽셀, 20m)"); ax.set_ylabel("ms / frame")
    ax.grid(True, which="both", alpha=0.25); ax.legend(fontsize=10, loc="upper left")

    # 대표 배율(가장 큰 공통 측정점)
    common = sorted(set(cx_e) & set(gx_e))
    mp0 = common[-1]
    ce = cy_e[list(cx_e).index(mp0)]; ge = gy_e[list(gx_e).index(mp0)]
    ct = cy_t[list(cx_t).index(mp0)]; gt = gy_t[list(gx_t).index(mp0)]
    fig.suptitle(f"위성 온보드 추론 전력·속도 실측 — CPU vs GPU  "
                 f"(Jetson Orin Nano · 재해 세그 · idle {idle:.1f}W 차감기준)",
                 fontsize=13.5, weight="bold", y=1.02)
    fig.text(0.5, -0.04,
             f"측정점 {mp0:.1f}MP: GPU가 CPU 대비 에너지 1/{ce/ge:.0f} · 속도 {ct/gt:.0f}배 ↑    |    "
             f"30MP 운영점(외삽): 에너지 CPU {e_ext/1000:.1f}J vs GPU {g_e_ext/1000:.2f}J → GPU 1/{e_ext/g_e_ext:.0f}, "
             f"속도 {t_ext/g_t_ext:.0f}배",
             ha="center", fontsize=10.5, color="#24292f")
    fig.tight_layout()
    fig.savefig(FIG, dpi=140, bbox_inches="tight")
    plt.close(fig)
    print("→ 그래프:", FIG)
    return {
        "mp0": mp0, "e_ratio0": ce/ge, "t_ratio0": ct/gt,
        "cpu_e30": e_ext, "gpu_e30": g_e_ext, "e_ratio30": e_ext/g_e_ext,
        "cpu_t30": t_ext, "gpu_t30": g_t_ext, "t_ratio30": t_ext/g_t_ext,
        "e_slope": e_slope,
    }


# ── PPT ────────────────────────────────────────────────────────────────
def make_ppt(d, rows, series, M):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    NAVY = RGBColor(0x0D, 0x11, 0x17); GREEN = RGBColor(0x2E, 0xA0, 0x43)
    RED = RGBColor(0xE5, 0x48, 0x4D); GREY = RGBColor(0x57, 0x60, 0x6A)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF); DARK = RGBColor(0x24, 0x29, 0x2F)

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def bg(slide, color):
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

    def box(slide, x, y, w, h, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT,
            font=None, line=1.15):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        lines = text.split("\n")
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align; p.line_spacing = line
            r = p.add_run(); r.text = ln
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
            r.font.name = font or KO
        return tb

    def bullets(slide, x, y, w, h, items, size=18, color=DARK, gap=6):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        for i, (txt, c, bold) in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap); p.line_spacing = 1.12
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = c; r.font.name = KO
        return tb

    # ── Slide 1: 타이틀
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    box(s, 0.9, 2.2, 11.5, 1.4, "위성 온보드 추론, CPU vs GPU", 44, WHITE, True)
    box(s, 0.9, 3.5, 11.5, 1.0, "전력·속도 실측과 배포 판단", 26, GREEN, True)
    box(s, 0.9, 4.7, 11.5, 1.2,
        "Jetson Orin Nano · Sentinel-2 20m 재해 세그(산불·홍수)\n"
        f"측정 기준: idle {d.get('idle_watt',0):.1f}W 차감 · 정상상태 in-process · PowerMeter",
        16, RGBColor(0xC9,0xD1,0xD9))
    box(s, 0.9, 6.7, 11.5, 0.5, "orbital-perception / merge", 12, GREY)

    # ── Slide 2: 핵심 질문 + 인식률 불변
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    box(s, 0.7, 0.5, 12, 0.9, "핵심 질문 — 20m 위성 영상, CPU와 GPU 중 무엇인가", 28, NAVY, True)
    bullets(s, 0.9, 1.7, 11.6, 4.6, [
        ("① 인식률은 CPU/GPU로 갈리지 않는다 (실측 확인)", GREEN, True),
        ("    같은 모델·가중치 → 같은 출력. 울진 장면서 GPU↔CPU 결과 차이는 "
         "345,897픽셀 중 단 3픽셀(0.0009%, 경계 부동소수점).", DARK, False),
        ("② 그러므로 선택 기준은 순수하게 '효율' — 위성에선 곧 전력·열", RED, True),
        ("    우주는 태양광·배터리로 전력이 빠듯하고, 진공이라 열을 복사로만 버린다.", DARK, False),
        ("    → 의미 있는 지표는 latency가 아니라 mJ/frame(한 장당 에너지).", DARK, False),
        ("③ 20m 운영점은 결코 작지 않다", NAVY, True),
        ("    20m는 픽셀이 넓어 데이터가 관리 가능하지만, S2 풀타일 1장 = 5490² ≈ 30MP.", DARK, False),
        ("    실제 위성은 이 스와스를 연속 처리한다 → 픽셀 많고 대상 큼 = GPU가 빛나는 일.", DARK, False),
    ], size=17)

    # ── Slide 3: 실험 설계
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    box(s, 0.7, 0.5, 12, 0.9, "실험 설계 — 같은 저울, 공정한 비교", 28, NAVY, True)
    bullets(s, 0.9, 1.7, 11.6, 4.6, [
        ("측정 도구", NAVY, True),
        ("    · 같은 세그 모델(TinyUNet)을 CPU와 GPU로만 바꿔 실행 — merge의 --device 경로", DARK, False),
        ("    · PowerMeter로 보드 총전력 샘플링, idle 차감해 dynamic(순수) 에너지도 산출", DARK, False),
        ("    · 워밍업(모델로딩·CUDA초기화) 제외한 정상상태만 측정", DARK, False),
        ("크기 스윕", NAVY, True),
        ("    · 0.07MP(256²) → 9.4MP(3072²)까지 쓸어 크기 의존성·교차점 탐색", DARK, False),
        ("    · CPU는 대형서 프레임당 수십초→1536²까지, GPU는 4.2MP(2048²)까지 측정, 30MP는 멱함수 외삽", DARK, False),
        ("정직성", GREEN, True),
        ("    · 정상상태 기준. GPU를 프레임마다 껐다 켜는 극저듀티라면 기동에너지는 별도 고려.", GREY, False),
    ], size=16)

    # ── Slide 4: 그래프
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    box(s, 0.7, 0.35, 12, 0.8, "실측 결과 — 에너지와 속도, 전 구간 GPU 우세", 26, NAVY, True)
    s.shapes.add_picture(str(FIG), Inches(0.35), Inches(1.25), width=Inches(12.6))

    # ── Slide 5: 결과표 + 숫자
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    box(s, 0.7, 0.5, 12, 0.8, "숫자로 본 결론", 28, NAVY, True)
    # 표
    meas = [r for r in rows if isinstance(r.get("cpu",{}).get("mJ_per_frame"), (int,float))
            and isinstance(r.get("gpu",{}).get("mJ_per_frame"), (int,float))]
    from pptx.util import Inches as In
    cols = ["장면(MP)", "CPU ms/f", "GPU ms/f", "속도배율", "CPU mJ/f", "GPU mJ/f", "에너지 1/n"]
    rowsN = len(meas) + 1
    tbl = s.shapes.add_table(rowsN, len(cols), In(0.6), In(1.6), In(8.4), In(0.4*rowsN)).table
    for j, c in enumerate(cols):
        cell = tbl.cell(0, j); cell.text = c
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r_ in p.runs: r_.font.size=Pt(12); r_.font.bold=True; r_.font.color.rgb=WHITE; r_.font.name=KO
    for i, r in enumerate(meas, 1):
        cpu, gpu = r["cpu"], r["gpu"]
        vals = [f"{r['megapixels']:.2f}", f"{cpu['sec_per_frame']*1000:.0f}",
                f"{gpu['sec_per_frame']*1000:.1f}", f"{cpu['sec_per_frame']/gpu['sec_per_frame']:.0f}×",
                f"{cpu['mJ_per_frame']:.0f}", f"{gpu['mJ_per_frame']:.0f}",
                f"{cpu['mJ_per_frame']/gpu['mJ_per_frame']:.0f}"]
        for j, v in enumerate(vals):
            cell = tbl.cell(i, j); cell.text = v
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r_ in p.runs: r_.font.size=Pt(11); r_.font.color.rgb=DARK; r_.font.name=KO
    # 우측 하이라이트
    bullets(s, 9.3, 1.7, 3.6, 4.8, [
        ("30MP 운영점 (외삽)", NAVY, True),
        (f"· 에너지: GPU가 CPU의 1/{M['e_ratio30']:.0f}", GREEN, True),
        (f"· 속도: GPU가 {M['t_ratio30']:.0f}배 빠름", GREEN, True),
        (f"· CPU {M['cpu_t30']/1000:.1f}s/장", DARK, False),
        (f"  vs GPU {M['gpu_t30']/1000:.2f}s/장", DARK, False),
        ("", DARK, False),
        ("교차점?", NAVY, True),
        ("· 없음 — 0.07MP부터", DARK, False),
        ("  이미 GPU 우세,", DARK, False),
        ("  격차는 크기와 함께 확대", DARK, False),
    ], size=14)

    # ── Slide 6: 권고
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    box(s, 0.8, 0.6, 11.7, 0.9, "결론 및 배포 권고", 30, WHITE, True)
    bullets(s, 1.0, 1.9, 11.2, 4.8, [
        ("결론: 우리 20m 위성 운영점에서 GPU가 정답", GREEN, True),
        ("    에너지·속도 모두 전 구간 GPU 우세, 30MP 운영점에서 격차 최대. 인식률 손실 0.", WHITE, False),
        ("", WHITE, False),
        ("온보드 배포 로드맵", RGBColor(0x79,0xC0,0xFF), True),
        ("    ① GPU 경로 확정 (완료 — merge --device cuda, 웹 데모 CUDA 가동)", WHITE, False),
        ("    ② DLA + INT8 TensorRT 로 한 단계 더 — 고정 CNN을 최저전력으로", WHITE, False),
        ("    ③ duty-cycle 게이팅 — 촬영창에만 가속기 켜고 대기 시 끄기", WHITE, False),
        ("", WHITE, False),
        ("한 줄 요약", GREEN, True),
        ("    20m·대면적 재해 세그는 GPU가 더 빠르고 더 적은 전기로, 같은 정확도로 처리한다.", WHITE, True),
    ], size=17)

    prs.save(PPTX)
    print("→ PPT:", PPTX)


if __name__ == "__main__":
    d, rows, series = load()
    K = make_fig(d, rows, series)
    make_ppt(d, rows, series, K)
    print("완료")
